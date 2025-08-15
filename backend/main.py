import os
import re
import shutil
import uuid
import json
import base64
import time
from datetime import datetime
from typing import List, Dict, Optional

import pandas as pd
import requests
import docx
import fitz
import pptx

from fastapi import FastAPI, UploadFile, File, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel

# ----------------- App Setup -----------------
app = FastAPI(
    title="Ollama LLM Recommender",
    version="33.0.0 (Final Integrated)",
    description="A robust, professionally architected AI Assistant with intelligent, multi-layered routing."
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ----------------- Constants & Setup -----------------
SERPER_API_KEY = "1e79c81fda1a17c6e0820ca2b5e04964a0573508"
UPLOAD_DIR = "uploads"
DATA_DIR = "data"
CHAT_HISTORY_FILE = os.path.join(DATA_DIR, "chat_history.json")
LLM_DATA_FILE = "llm_data.csv"
VISION_MODEL = "llava"
TEXT_MODEL = "phi3"

SYSTEM_PROMPT = """You are a helpful and friendly AI Assistant from Caze Labs. 

Your primary capabilities include:
- Recommending LLM models based on user needs
- Processing and analyzing documents (PDF, Word, PowerPoint, text files)
- Analyzing and describing images
- Answering general questions with web search when needed

Always respond in a natural, conversational, and professional way. Be helpful and engaging while staying focused on your core capabilities."""

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(DATA_DIR, exist_ok=True)

# ----------------- Data Loading & Persistence -----------------
try:
    llm_df = pd.read_csv(LLM_DATA_FILE)
    print(f"✅ Loaded {len(llm_df)} LLM models from database")
except FileNotFoundError:
    print(f"⚠️  WARNING: '{LLM_DATA_FILE}' not found. LLM recommendations will not work.")
    llm_df = pd.DataFrame()

def load_json_data(filepath: str) -> Dict:
    if os.path.exists(filepath):
        try:
            with open(filepath, 'r', encoding='utf-8') as f: return json.load(f)
        except (json.JSONDecodeError, IOError): return {}
    return {}

def save_json_data(data: Dict, filepath: str):
    try:
        with open(filepath, 'w', encoding='utf-8') as f: json.dump(data, f, indent=4)
    except IOError as e: print(f"Error saving data: {e}")

# ----------------- Core AI and Tool Functions -----------------

def get_ollama_response(prompt: str, model: str = TEXT_MODEL, images: list = None, retries: int = 2, delay: int = 3) -> str:
    payload = {"model": model, "prompt": prompt, "stream": False}
    if images:
        payload["images"] = images
    for attempt in range(retries):
        try:
            response = requests.post("http://localhost:11434/api/generate", json=payload, timeout=180)
            response.raise_for_status()
            return response.json().get("response", "").strip()
        except requests.exceptions.RequestException as e:
            print(f"Ollama connection error (Attempt {attempt + 1}/{retries}): {e}")
            if attempt < retries - 1:
                time.sleep(delay)
            else:
                return "⚠️ I'm having trouble connecting to the AI model right now. Please try again in a moment."
    return "⚠️ The AI model failed to respond after several attempts."

# --- Document Processing Functions ---
def read_text_safely(file_path: str) -> str:
    try:
        with open(file_path, 'r', encoding='utf-8') as f: return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='latin-1') as f: return f.read()

def read_pdf(file_path: str) -> str:
    with fitz.open(file_path) as doc:
        return "".join(page.get_text() for page in doc)

def read_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def read_pptx(file_path: str) -> str:
    prs = pptx.Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    return "\n".join(text_runs)

def process_document(file_path: str, user_prompt: str) -> str:
    content = ""
    try:
        file_extension = os.path.splitext(file_path)[1].lower()
        if file_extension in ['.txt', '.md', '.py', '.json', '.csv']:
            content = read_text_safely(file_path)
        elif file_extension == '.pdf':
            content = read_pdf(file_path)
        elif file_extension == '.docx':
            content = read_docx(file_path)
        elif file_extension == '.pptx':
            content = read_pptx(file_path)
        else:
            return f"Unsupported file type: '{file_extension}'. I can currently read PDF, DOCX, PPTX, and plain text files."
        
        prompt_template = f"Based on the following document content, please answer the user's question comprehensively.\n\nDOCUMENT CONTENT:\n---\n{content[:5000]}\n---\n\nUSER QUESTION: {user_prompt}"
        return get_ollama_response(prompt_template)
    except Exception as e:
        return f"I'm sorry, I encountered an error while processing the file: {str(e)}"

def describe_image(image_path: str, prompt: str) -> str:
    try:
        with open(image_path, "rb") as image_file:
            encoded_string = base64.b64encode(image_file.read()).decode('utf-8')
        vision_prompt = f"Please analyze this image and respond to the user's request: {prompt}"
        return get_ollama_response(vision_prompt, model=VISION_MODEL, images=[encoded_string])
    except Exception as e:
        return f"Error processing image: {str(e)}"

def search_local_models(prompt: str) -> str:
    if llm_df.empty:
        return "I cannot search for models because my database file (llm_data.csv) is missing."
    
    keywords = [word for word in prompt.lower().split() if len(word) > 2]
    tags_pattern = '|'.join(keywords)
    search_pattern = f'(?i)({tags_pattern})'
    recommended_models = llm_df[llm_df['use_case_tags'].str.contains(search_pattern, na=False, regex=True)]
    
    if not recommended_models.empty:
        bot_reply = f"Based on your request, here are my top LLM recommendations:\n\n"
        for i, (_, model) in enumerate(recommended_models.head(3).iterrows()):
            bot_reply += f"### {i+1}. **{model['model_name']}**\n"
            bot_reply += f"- **Primary Use Cases:** {model['use_case_tags']}\n"
            bot_reply += f"- **Size (GB):** {model['size_gb']}\n"
            bot_reply += f'- **Link:** <a href="{model["link"]}" target="_blank" style="color: #60a5fa;">{model["source"]} Page</a>\n'
            bot_reply += f"- **Notes:** {model['notes']}\n\n"
        bot_reply += "Need more specific recommendations? Feel free to tell me more about your exact use case!"
        return bot_reply
    else:
        return f"I searched my database but couldn't find models matching your specific criteria. Try using keywords like 'coding', 'chat', 'writing', or 'analysis' for better results."

def perform_web_search(query: str) -> str:
    headers = {'X-API-KEY': SERPER_API_KEY, 'Content-Type': 'application/json'}
    payload = json.dumps({"q": query})
    try:
        response = requests.post("https://google.serper.dev/search", headers=headers, data=payload, timeout=10)
        response.raise_for_status()
        results = response.json()
        if "answerBox" in results:
            return f"Here's what I found: {results['answerBox']['answer']}"
        if "organic" in results and results["organic"]:
            snippets = [f"**{r.get('title', 'No title')}**: {r.get('snippet', 'No description available')}" for r in results["organic"][:3]]
            return "Here are the top search results:\n\n" + "\n\n".join(snippets)
        return "I couldn't find relevant information for that query. Could you try rephrasing your question?"
    except Exception as e:
        print(f"Web search error: {e}")
        return f"I encountered an issue while searching the web. Let me try to help you with my existing knowledge instead."

# --- NEW HELPER FUNCTIONS FOR INTELLIGENT ROUTING (FROM YOUR DESIGN) ---

def is_llm_recommendation_request(prompt_lower: str) -> bool:
    strong_indicators = ["recommend llm", "suggest llm", "best llm", "which llm", "recommend model", "suggest model", "best model", "which model", "llm for", "model for"]
    if any(indicator in prompt_lower for indicator in strong_indicators):
        return True
    recommendation_words = ["recommend", "suggest", "best", "top", "good"]
    model_words = ["llm", "model", "language model"]
    has_recommendation = any(word in prompt_lower for word in recommendation_words)
    has_model = any(word in prompt_lower for word in model_words)
    if has_recommendation and has_model:
        general_questions = ["what is llm", "what are llm", "explain llm", "define llm"]
        if not any(question in prompt_lower for question in general_questions):
            return True
    return False

def handle_ai_education_question(prompt: str) -> str:
    prompt_lower = prompt.lower()
    if "llm" in prompt_lower:
        return """**LLM** stands for **Large Language Model** - these are AI systems trained on vast amounts of text data to understand and generate human-like text.
**Key characteristics:**
• **Large Scale**: Trained on billions of parameters and massive datasets
• **Versatile**: Can handle multiple tasks like writing, coding, analysis, conversation
• **Generative**: Create new content rather than just retrieving information
Would you like me to recommend specific LLMs for particular use cases?"""
    elif any(ai_term in prompt_lower for ai_term in ["what is ai", "artificial intelligence", "how does ai work"]):
        return """**Artificial Intelligence (AI)** is technology that enables machines to simulate human intelligence.
**Key AI capabilities**:
• **Learning**: Improving performance through experience
• **Reasoning**: Making logical connections and decisions
• **Problem-solving**: Finding solutions to complex challenges
I'm here to demonstrate AI in action! What would you like to explore?"""
    return get_ollama_response(f"Please explain this AI/technology concept in a conversational, educational way: {prompt}")

def should_use_web_search(prompt_lower: str) -> bool:
    web_search_indicators = ["current", "latest", "recent", "today", "now", "news", "update", "status", "trending"]
    if any(indicator in prompt_lower for indicator in web_search_indicators):
        return True
    complex_indicators = ["explain", "tell me about", "what is", "how does", "why does"]
    if len(prompt_lower.split()) > 4 and any(indicator in prompt_lower for indicator in complex_indicators):
        return True
    return False

def handle_natural_conversation(prompt: str, prompt_lower: str) -> str:
    if any(identity in prompt_lower for identity in ["who are you", "what are you"]):
        return """Hello! I'm an AI Assistant from Caze Labs. I'm designed to be helpful, conversational, and intelligent. I can help you with:
• **Finding AI Models**
• **Document Analysis**
• **Image Processing**
• **General Questions**
What would you like to explore?"""
    elif any(greeting in prompt_lower for greeting in ["hi", "hello", "hey"]):
        return "Hello! How can I help you today?"
    elif any(thanks in prompt_lower for thanks in ["thank", "thanks"]):
        return "You're very welcome! Is there anything else I can help with?"
    elif any(goodbye in prompt_lower for goodbye in ["bye", "goodbye"]):
        return "Goodbye! Have a great day."
    elif any(concept in prompt_lower for concept in ["what is llm", "what is ai"]):
        return handle_ai_education_question(prompt)
    elif should_use_web_search(prompt_lower):
        print("🌐 Using web search for comprehensive answer")
        search_result = perform_web_search(prompt)
        enhanced_prompt = f"""You are an intelligent AI assistant. A user asked: "{prompt}"
Here's information from a web search: {search_result}
Please provide a comprehensive, conversational, and helpful response using this information."""
        return get_ollama_response(enhanced_prompt)
    else:
        print("🤖 Direct AI conversation")
        conversational_prompt = f"""You are a helpful, intelligent AI assistant from Caze Labs.
User: {prompt}
Respond in a natural, conversational way."""
        return get_ollama_response(conversational_prompt)

# ----------------- Pydantic Models & Main Endpoint -----------------
class ChatMessage(BaseModel): id: str; role: str; text: Optional[str] = None; file_url: Optional[str] = None; timestamp: str

@app.post("/chat/", response_model=List[ChatMessage])
async def chat_recommender(prompt: str = Form(...), file: Optional[UploadFile] = File(None), user_id: str = Form("default_user")):
    chat_history = load_json_data(CHAT_HISTORY_FILE)
    
    file_path, file_url = None, None
    if file:
        file_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
        with open(file_path, "wb") as buffer: shutil.copyfileobj(file.file, buffer)
        file_url = f"/uploads/{os.path.basename(file_path)}"

    user_msg = ChatMessage(id=str(uuid.uuid4()), role="user", text=prompt, file_url=file_url, timestamp=datetime.now().isoformat())
    chat_history.setdefault(user_id, []).append(user_msg.dict())

    final_response = ""
    prompt_lower = prompt.lower().strip()
    
    try:
        if file_path:
            file_extension = os.path.splitext(file_path)[1].lower()
            if file_extension in ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp']:
                print("🖼️  Routing to: Image Analysis")
                final_response = describe_image(file_path, prompt)
            else:
                print("📄 Routing to: Document Processing")
                final_response = process_document(file_path, prompt)
        elif is_llm_recommendation_request(prompt_lower):
            print("🔍 Routing to: LLM Model Search")
            final_response = search_local_models(prompt)
        else:
            print("💬 Routing to: Natural Conversation")
            final_response = handle_natural_conversation(prompt, prompt_lower)
    except Exception as e:
        print(f"❌ Error during routing: {e}")
        final_response = "I encountered an error while processing your request. Please try again!"

    if not final_response or final_response.strip() == "":
        final_response = "I'm here to help! Feel free to ask me anything."

    bot_msg = ChatMessage(id=str(uuid.uuid4()), role="bot", text=final_response, timestamp=datetime.now().isoformat())
    chat_history.setdefault(user_id, []).append(bot_msg.dict())
    save_json_data(chat_history, CHAT_HISTORY_FILE)
    
    return chat_history.get(user_id, [])

# ----------------- Additional Endpoints -----------------
@app.get("/")
def root():
    return {"status": "✅ Caze Labs AI Assistant is running", "version": "33.1.0"}

# ... (other endpoints like /health, etc. can be added here)
